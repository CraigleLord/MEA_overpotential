"""
PPT generator for:
"Thermostable Tellurium Anchoring Enabling Robust Thermal and Electrochemical
 Stability for Pt3Co Intermetallic Fuel Cell Catalysts"
Chen et al., Advanced Functional Materials, 2024
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── paths ──────────────────────────────────────────────────────────────────
BASE = r"c:\Users\user\My Drive\KAIST MASc 2021\Laboratory Work\Protocol\overpotential calculation\For paper SI"
FIGS = os.path.join(BASE, "paper_figs")
OUT  = os.path.join(BASE, "Te_Pt3Co_PEMFC_talk.pptx")

# ── colours ────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0D, 0x2B, 0x55)   # dark navy (slide backgrounds / headers)
BLUE   = RGBColor(0x1A, 0x5C, 0x9A)   # accent blue
TEAL   = RGBColor(0x00, 0x7A, 0x87)   # accent teal
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF4, 0xF6, 0xF8)   # very light grey (body bg)
DKGRAY = RGBColor(0x2C, 0x3E, 0x50)
RED    = RGBColor(0xC0, 0x39, 0x2B)

# ── helpers ────────────────────────────────────────────────────────────────
W, H = Inches(13.33), Inches(7.5)      # widescreen 16:9

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # truly blank
    return prs.slides.add_slide(blank_layout)

def fill_bg(slide, color):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def txb(slide, text, left, top, width, height,
        font_size=20, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def add_image(slide, path, left, top, width=None, height=None):
    if width and height:
        slide.shapes.add_picture(path, left, top, width, height)
    elif width:
        slide.shapes.add_picture(path, left, top, width=width)
    elif height:
        slide.shapes.add_picture(path, left, top, height=height)
    else:
        slide.shapes.add_picture(path, left, top)

def header_bar(slide, title_text, subtitle=None):
    """Navy bar across top with white title."""
    add_rect(slide, 0, 0, W, Inches(1.1), fill_color=NAVY)
    txb(slide, title_text,
        Inches(0.25), Inches(0.08), Inches(12.5), Inches(0.75),
        font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txb(slide, subtitle,
            Inches(0.25), Inches(0.78), Inches(12.5), Inches(0.35),
            font_size=14, bold=False, color=RGBColor(0xAA,0xCC,0xFF),
            align=PP_ALIGN.LEFT)

def bullet_box(slide, items, left, top, width, height,
               font_size=17, color=DKGRAY, bullet="▸ "):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = bullet + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color

def key_number(slide, value, label, left, top, width=Inches(2.4), height=Inches(1.2),
               bg=BLUE, val_size=32, lbl_size=13):
    add_rect(slide, left, top, width, height, fill_color=bg)
    txb(slide, value, left, top+Inches(0.08), width, Inches(0.65),
        font_size=val_size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(slide, label, left, top+Inches(0.7), width, Inches(0.45),
        font_size=lbl_size, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════

prs = new_prs()

# ─── SLIDE 1 – Title ──────────────────────────────────────────────────────────
sl = blank_slide(prs)
fill_bg(sl, NAVY)

# Decorative teal accent bar (left)
add_rect(sl, 0, 0, Inches(0.18), H, fill_color=TEAL)

# Journal badge
add_rect(sl, Inches(0.35), Inches(0.3), Inches(3.2), Inches(0.42), fill_color=BLUE)
txb(sl, "Advanced Functional Materials  |  2024",
    Inches(0.35), Inches(0.3), Inches(3.6), Inches(0.42),
    font_size=12, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Title
txb(sl,
    "Thermostable Tellurium Anchoring Enabling Robust\n"
    "Thermal and Electrochemical Stability for\n"
    "Pt₃Co Intermetallic Fuel Cell Catalysts",
    Inches(0.35), Inches(1.0), Inches(12.5), Inches(2.8),
    font_size=38, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Authors
txb(sl,
    "Yuanxin Chen, Zihan Meng, Fei Liu, Aojie Zhang, Xiaocan Wang,\n"
    "Yifei Xiong, Haibo Tang, Tian Tian, Haolin Tang",
    Inches(0.35), Inches(3.8), Inches(12.0), Inches(0.9),
    font_size=16, bold=False, color=RGBColor(0xAA,0xCC,0xFF), align=PP_ALIGN.LEFT)

# Presenter line
add_rect(sl, Inches(0.35), Inches(4.8), Inches(4.5), Inches(0.04), fill_color=TEAL)
txb(sl, "Literature Talk Presentation  |  2026",
    Inches(0.35), Inches(4.9), Inches(6.0), Inches(0.4),
    font_size=14, bold=False, color=RGBColor(0xCC,0xDD,0xEE))

# ─── SLIDE 2 – Motivation / Background ───────────────────────────────────────
sl = blank_slide(prs)
fill_bg(sl, LGRAY)
header_bar(sl, "Background: Why Better PEMFC Catalysts?")

# Left column – text
bullet_box(sl, [
    "PEMFCs convert H₂ → electricity with zero emissions; only water byproduct",
    "Sluggish oxygen reduction reaction (ORR) at the cathode limits performance",
    "Pt nanoparticles on carbon are state-of-the-art, but: high cost, low durability",
    "Alloying Pt with 3d metals (Co, Ni, Mn) boosts activity via electronic / strain effects",
    "Ordered intermetallic Pt₃Co (L1₂ structure) = thermodynamically stable, high activity",
    "Key challenge: synthesising ordered i-NPs requires high-T annealing ➜ sintering",
    "Carbon supports often have weak metal-support interactions (MSIs) ➜ particle growth",
], Inches(0.3), Inches(1.25), Inches(6.5), Inches(5.8),
   font_size=17, color=DKGRAY)

# Right column – key numbers box
add_rect(sl, Inches(7.0), Inches(1.3), Inches(5.9), Inches(5.8), fill_color=NAVY)
txb(sl, "DOE 2025 Targets", Inches(7.1), Inches(1.4), Inches(5.7), Inches(0.5),
    font_size=18, bold=True, color=WHITE)
key_number(sl, "≥ 0.44 A/mgPt", "Mass Activity @ 0.9 V",
           Inches(7.2), Inches(2.0), width=Inches(2.5), height=Inches(1.3), bg=BLUE)
key_number(sl, "≥ 60%", "MA Retention after 30k cycles",
           Inches(10.0), Inches(2.0), width=Inches(2.5), height=Inches(1.3), bg=BLUE)
key_number(sl, "≤ 0.2 mgPt/cm²", "Cathode Pt Loading",
           Inches(7.2), Inches(3.5), width=Inches(2.5), height=Inches(1.3), bg=TEAL)
key_number(sl, "< 5 nm", "Target NP size for ECSA",
           Inches(10.0), Inches(3.5), width=Inches(2.5), height=Inches(1.3), bg=TEAL)
txb(sl, "This work targets all four simultaneously.",
    Inches(7.2), Inches(5.1), Inches(5.5), Inches(0.6),
    font_size=15, bold=True, color=RGBColor(0x88,0xFF,0xCC))

# ─── SLIDE 3 – The Problem ────────────────────────────────────────────────────
sl = blank_slide(prs)
fill_bg(sl, LGRAY)
header_bar(sl, "The Problem: Sintering During High-Temperature Annealing",
           subtitle="Ordering Pt₃Co i-NPs requires ≥700°C — but this causes severe particle growth on bare carbon")

bullet_box(sl, [
    "High-T annealing is necessary to drive atomic diffusion and achieve the ordered L1₂ Pt₃Co phase",
    "On bare carbon (C), NPs grow to ~9.1 nm — far too large for practical ECSA",
    "Two sintering pathways: particle migration & coalescence (PMC) + Ostwald ripening (OR)",
    "Weak Pt–C interactions allow NP detachment under fuel cell operating conditions",
    "Previous solutions: S-doped C, P-O groups on Ketjen Black — partial success only",
    "Group VIA dopants (S, Se, Te) can enhance MSIs by modulating electron transfer",
    "Te has the lowest electronegativity among S/Se/Te → stronger metallicity, higher conductivity (~1000 S/m), superior acid resistance",
], Inches(0.3), Inches(1.25), Inches(7.8), Inches(5.5), font_size=17, color=DKGRAY)

# Insight box
add_rect(sl, Inches(8.2), Inches(1.3), Inches(4.8), Inches(5.8), fill_color=NAVY)
txb(sl, "Key Insight", Inches(8.3), Inches(1.4), Inches(4.5), Inches(0.45),
    font_size=18, bold=True, color=TEAL)
txb(sl,
    "Te is thermostable at high temperatures — "
    "while S and Se are volatilised during annealing, "
    "Te remains on the carbon surface at 2.3 wt%, "
    "maintaining strong MSIs throughout synthesis AND operation.",
    Inches(8.3), Inches(2.0), Inches(4.4), Inches(3.5),
    font_size=16, color=WHITE)
txb(sl, "→ This is the core novelty of the paper.",
    Inches(8.3), Inches(5.5), Inches(4.4), Inches(0.5),
    font_size=15, bold=True, color=RGBColor(0x88,0xFF,0xCC))

# ─── SLIDE 4 – Strategy / Overview ───────────────────────────────────────────
sl = blank_slide(prs)
fill_bg(sl, LGRAY)
header_bar(sl, "Strategy: Te-C Support via Chemical Vapour Deposition",
           subtitle="Figure 1a — Schematic of carbon modification and Pt₃Co i-NP synthesis")

add_image(sl, os.path.join(FIGS, "fig1.png"),
          Inches(0.2), Inches(1.15), width=Inches(13.0))

txb(sl,
    "CVD deposits Te onto Ketjen Black (EC-300J) → Te–C support    "
    "→ Pt/Co precursors impregnated → high-T anneal (alloying) + low-T hold (ordering) "
    "→ Pt₃Co i-NPs with L1₂ structure anchored on Te–C",
    Inches(0.2), Inches(6.9), Inches(13.0), Inches(0.55),
    font_size=14, color=DKGRAY, align=PP_ALIGN.CENTER)

# ─── SLIDE 5 – Structural Characterisation: Particle Size ────────────────────
sl = blank_slide(prs)
fill_bg(sl, LGRAY)
header_bar(sl, "Structural Characterisation: Particle Size & Ordering",
           subtitle="Figures 1b–f — XRD and HAADF-STEM")

add_image(sl, os.path.join(FIGS, "fig1.png"),
          Inches(0.2), Inches(1.15), width=Inches(8.5))

# Stats table on right
add_rect(sl, Inches(8.9), Inches(1.2), Inches(4.2), Inches(5.5), fill_color=NAVY)
txb(sl, "Particle Size Comparison", Inches(9.0), Inches(1.3), Inches(4.0), Inches(0.45),
    font_size=16, bold=True, color=WHITE)

rows = [
    ("Catalyst",       "Size (XRD)",  "I₁₁₀/I₁₁₁"),
    ("Pt₃Co/Te-C",    "3.9 nm",      "0.212"),
    ("Pt₃Co/Se-C",   "5.6 nm",      "0.137"),
    ("Pt₃Co/S-C",    "7.8 nm",      "0.113"),
    ("Pt₃Co/C",      "9.1 nm",      "0.105"),
]
for ri, row in enumerate(rows):
    bg = BLUE if ri == 0 else (RGBColor(0x1E,0x3A,0x5F) if ri % 2 else RGBColor(0x16,0x2D,0x4A))
    if ri == 1:
        bg = TEAL
    add_rect(sl, Inches(9.0), Inches(1.85)+ri*Inches(0.68),
             Inches(3.9), Inches(0.65), fill_color=bg)
    txb(sl, f"{row[0]}    {row[1]}    {row[2]}",
        Inches(9.05), Inches(1.9)+ri*Inches(0.68), Inches(3.8), Inches(0.6),
        font_size=14 if ri > 0 else 13, bold=(ri==0), color=WHITE)

txb(sl,
    "Te doping gives smallest NPs AND highest ordering.\n"
    "I₁₁₀/I₁₁₁ > standard L1₂ value → fully ordered phase.",
    Inches(9.0), Inches(5.3), Inches(4.0), Inches(1.3),
    font_size=14, color=RGBColor(0x88,0xFF,0xCC))

# ─── SLIDE 6 – Atomic-Resolution STEM ───────────────────────────────────────
sl = blank_slide(prs)
fill_bg(sl, LGRAY)
header_bar(sl, "Atomic-Resolution Analysis: L1₂ Ordered Structure Confirmed",
           subtitle="Figures 1g–j — HAADF-STEM, FFT, line profile, EDS mapping of Pt₃Co/Te-C")

add_image(sl, os.path.join(FIGS, "fig1.png"),
          Inches(0.2), Inches(1.15), width=Inches(8.5))

bullet_box(sl, [
    "4 nm Pt₃Co i-NP imaged along [001] direction",
    "HAADF Z-contrast: Pt columns brighter than Co",
    "FFT patterns confirm L1₂ ordered superlattice",
    "Alternating intensity profiles verify atomic ordering",
    "EDS mapping: Pt and Co uniformly distributed; Pt/Co ≈ 3:1 (ideal)",
    "Residual Te enriched near Pt₃Co i-NPs → evidence of Pt–Te bonding",
], Inches(9.0), Inches(1.4), Inches(4.1), Inches(5.5),
   font_size=16, color=DKGRAY)

# ─── SLIDE 7 – In-situ Thermal Stability ─────────────────────────────────────
sl = blank_slide(prs)
fill_bg(sl, LGRAY)
header_bar(sl, "In-situ Thermal Stability: Te Suppresses Sintering During Annealing",
           subtitle="Figure 2 — In-situ high-temperature XRD")

add_image(sl, os.path.join(FIGS, "fig2.png"),
          Inches(0.3), Inches(1.2), width=Inches(8.0))

bullet_box(sl, [
    "In-situ XRD tracks (111) peak width vs. temperature up to 700 °C",
    "Pt₃Co/Te-C: particle size grows to only 4.13 nm after 6 h at 700°C",
    "Pt₃Co/C: size reaches ~7 nm under the same conditions",
    "Superlattice (110) peak clearly visible in Pt₃Co/Te-C → ordering in progress",
    "(110) peak absent in Pt₃Co/C → Te is essential for ordering",
    "Te provides a thermal anchor preventing PMC and Ostwald ripening",
], Inches(8.5), Inches(1.3), Inches(4.6), Inches(5.5),
   font_size=16, color=DKGRAY)

# ─── SLIDE 8 – Electronic Structure (XPS) ────────────────────────────────────
sl = blank_slide(prs)
fill_bg(sl, LGRAY)
header_bar(sl, "Electronic Structure: Te Doping Strengthens Metal–Support Interactions",
           subtitle="Figure 3a — XPS Pt 4f spectra")

add_image(sl, os.path.join(FIGS, "fig3.png"),
          Inches(0.3), Inches(1.2), width=Inches(7.8))

add_rect(sl, Inches(8.3), Inches(1.2), Inches(4.8), Inches(5.8), fill_color=NAVY)
txb(sl, "Binding Energy Shifts (vs Pt₃Co/C)",
    Inches(8.4), Inches(1.3), Inches(4.6), Inches(0.5),
    font_size=16, bold=True, color=WHITE)

shifts = [
    ("Pt 4f₇/₂ in Pt₃Co/Te-C", "+0.120 eV", TEAL),
    ("Pt 4f₇/₂ in Pt₃Co/Se-C", "+0.075 eV", BLUE),
    ("Pt 4f₇/₂ in Pt₃Co/S-C",  "+0.060 eV", RGBColor(0x5A,0x7A,0x9A)),
    ("Co 2p₁/₂ in Pt₃Co/Te-C", "+0.95 eV",  TEAL),
]
for i, (label, val, col) in enumerate(shifts):
    add_rect(sl, Inches(8.4), Inches(1.95)+i*Inches(0.85),
             Inches(4.5), Inches(0.75), fill_color=col)
    txb(sl, f"{label}\n→ {val}",
        Inches(8.5), Inches(2.0)+i*Inches(0.85), Inches(4.2), Inches(0.7),
        font_size=14, bold=False, color=WHITE)

txb(sl,
    "Positive BE shift = electron transfer from Pt/Co to Te-C\n"
    "→ d-band centre downshift → weaker O-species adsorption\n"
    "→ Enhanced ORR activity",
    Inches(8.4), Inches(5.4), Inches(4.5), Inches(1.2),
    font_size=14, color=RGBColor(0x88,0xFF,0xCC))

# ─── SLIDE 9 – XAS / EXAFS ───────────────────────────────────────────────────
sl = blank_slide(prs)
fill_bg(sl, LGRAY)
header_bar(sl, "Coordination Structure: Pt–Te Bonds Confirmed by XAS",
           subtitle="Figures 3b–e — XANES, EXAFS, wavelet transform, fitting")

add_image(sl, os.path.join(FIGS, "fig3.png"),
          Inches(0.3), Inches(1.2), width=Inches(7.8))

bullet_box(sl, [
    "XANES: Pt in Pt₃Co/Te-C ≈ metallic state (valence +0.13); similar to Pt foil white-line",
    "EXAFS R-space: broad peak at ~2.21 Å = overlapping Pt–Pt, Pt–Co, Pt–Te shells",
    "EXAFS fitting coordination numbers:",
    "   • Pt–Pt: CN = 7.7  (bond 2.60 Å)",
    "   • Pt–Co: CN = 2.2  (bond 2.68 Å)",
    "   • Pt–Te: CN = 1.3  (bond 2.65 Å) ← new bond from Te doping",
    "Compressive strain confirmed: Pt–Pt shorter than Pt foil (2.425 Å)",
    "Wavelet transform analysis corroborates overlapping multi-shell coordination",
], Inches(8.3), Inches(1.3), Inches(4.8), Inches(5.8),
   font_size=15, color=DKGRAY)

# ─── SLIDE 10 – Te Thermostability ──────────────────────────────────────────
sl = blank_slide(prs)
fill_bg(sl, LGRAY)
header_bar(sl, "Thermostability: Only Te Survives High-Temperature Annealing",
           subtitle="XPS + ICP-OES comparison before and after annealing")

# Large comparison table
add_rect(sl, Inches(0.3), Inches(1.3), Inches(12.7), Inches(4.8), fill_color=NAVY)
txb(sl, "Heteroatom Retention After High-Temperature Annealing",
    Inches(0.4), Inches(1.35), Inches(12.0), Inches(0.5),
    font_size=17, bold=True, color=WHITE)

cols = [("Heteroatom", ""), ("XPS Signal\nafter anneal", ""), ("ICP Retention\n(wt%)", ""), ("MSI Status", "")]
data = [
    ("Te", "Clearly visible (Te 3d)", "2.3 wt%", "Strong – Pt–Te bonds maintained"),
    ("Se", "Barely detectable",        "0.9 wt%", "Partial"),
    ("S",  "Not detected",             "0.2 wt% (nearly zero)", "Negligible"),
]
for ci, (hd, _) in enumerate(cols):
    add_rect(sl, Inches(0.4)+ci*Inches(3.1), Inches(1.95), Inches(3.0), Inches(0.55), fill_color=BLUE)
    txb(sl, hd, Inches(0.5)+ci*Inches(3.1), Inches(1.97), Inches(2.9), Inches(0.55),
        font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
for ri, row in enumerate(data):
    bg_row = TEAL if ri==0 else (RGBColor(0x1E,0x3A,0x5F) if ri%2==1 else RGBColor(0x14,0x2A,0x44))
    for ci, cell in enumerate(row):
        add_rect(sl, Inches(0.4)+ci*Inches(3.1), Inches(2.6)+ri*Inches(0.75),
                 Inches(3.0), Inches(0.7), fill_color=bg_row)
        txb(sl, cell, Inches(0.5)+ci*Inches(3.1), Inches(2.63)+ri*Inches(0.75),
            Inches(2.9), Inches(0.65), font_size=14, color=WHITE, align=PP_ALIGN.CENTER)

txb(sl,
    "Conclusion: Te's inherent thermostability ensures sufficient Te remains on the support "
    "to maintain strong MSIs not only during synthesis but throughout fuel cell operation.",
    Inches(0.3), Inches(6.2), Inches(12.7), Inches(0.7),
    font_size=16, bold=True, color=DKGRAY)

# ─── SLIDE 11 – ORR Activity (RDE) ──────────────────────────────────────────
sl = blank_slide(prs)
fill_bg(sl, LGRAY)
header_bar(sl, "ORR Electrochemical Activity: RDE Results",
           subtitle="Figure 4a–c — LSV, specific/mass activity, ECSA")

add_image(sl, os.path.join(FIGS, "fig4.png"),
          Inches(0.3), Inches(1.15), width=Inches(8.2))

add_rect(sl, Inches(8.6), Inches(1.2), Inches(4.5), Inches(5.8), fill_color=NAVY)
txb(sl, "Activity @ 0.9 V (0.1 M HClO₄)",
    Inches(8.7), Inches(1.3), Inches(4.3), Inches(0.45),
    font_size=16, bold=True, color=WHITE)

metrics = [
    ("E₁/₂",  "Pt₃Co/Te-C: 0.925 V", "Pt₃Co/C: 0.895 V"),
    ("SA",     "0.64 mA/cmPt²",       "0.41 mA/cmPt²"),
    ("MA",     "301 mA/mgPt",         "103 mA/mgPt"),
    ("ECSA",   "47.3 m²/gPt",         "25.2 m²/gPt"),
]
for i, (metric, te_val, c_val) in enumerate(metrics):
    y = Inches(1.95) + i*Inches(1.0)
    add_rect(sl, Inches(8.7), y, Inches(4.2), Inches(0.9),
             fill_color=TEAL if i==0 else RGBColor(0x1E,0x3A,0x5F))
    txb(sl, f"{metric}:  Te-C → {te_val}  |  C → {c_val}",
        Inches(8.8), y+Inches(0.05), Inches(4.0), Inches(0.8),
        font_size=13, color=WHITE)

txb(sl,
    "Pt₃Co/Te-C SA = 1.56× and MA = 2.92× higher than Pt₃Co/C\n"
    "Smaller NPs + stronger MSI = better ECSA and activity",
    Inches(8.7), Inches(5.9), Inches(4.2), Inches(0.9),
    font_size=14, bold=True, color=RGBColor(0x88,0xFF,0xCC))

# ─── SLIDE 12 – ORR Durability (RDE ADT) ────────────────────────────────────
sl = blank_slide(prs)
fill_bg(sl, LGRAY)
header_bar(sl, "Exceptional Electrochemical Durability: 100,000 Voltage Cycles",
           subtitle="Figure 4d–j — Accelerated durability test (ADT), 0.6–1.0 V, 0.1 M HClO₄")

add_image(sl, os.path.join(FIGS, "fig4.png"),
          Inches(0.3), Inches(1.15), width=Inches(8.2))

add_rect(sl, Inches(8.6), Inches(1.2), Inches(4.5), Inches(5.8), fill_color=NAVY)
txb(sl, "After ADT", Inches(8.7), Inches(1.3), Inches(4.2), Inches(0.45),
    font_size=16, bold=True, color=WHITE)

adt_data = [
    ("",             "Pt₃Co/Te-C\n(100k cycles)", "Pt₃Co/C\n(30k cycles)"),
    ("ΔE₁/₂",       "−8 mV",                     "−20 mV"),
    ("ΔECSA",        "−29.8%",                    "−16.3%"),
    ("ΔSA",          "−1.5%  ✓",                  "−5%"),
    ("ΔMA",          "−30%",                      "−20%"),
    ("Morphology",   "Unchanged ✓",               "Agglomerated ✗"),
]
for ri, row in enumerate(adt_data):
    bg = BLUE if ri==0 else (TEAL if ri==3 else RGBColor(0x1E,0x3A,0x5F) if ri%2==0 else RGBColor(0x14,0x2A,0x44))
    y = Inches(1.9) + ri*Inches(0.72)
    add_rect(sl, Inches(8.7), y, Inches(4.2), Inches(0.68), fill_color=bg)
    txb(sl, f"{row[0]:12s}  {row[1]:25s}  {row[2]}",
        Inches(8.75), y+Inches(0.03), Inches(4.1), Inches(0.63),
        font_size=12, color=WHITE)

txb(sl, "SA loss of only 1.5% after 100k cycles is unprecedented.",
    Inches(8.7), Inches(6.2), Inches(4.2), Inches(0.5),
    font_size=14, bold=True, color=RGBColor(0x88,0xFF,0xCC))

# ─── SLIDE 13 – MEA Performance ─────────────────────────────────────────────
sl = blank_slide(prs)
fill_bg(sl, LGRAY)
header_bar(sl, "MEA Performance: Exceeds DOE 2025 Targets",
           subtitle="Figure 5a–c — H₂–O₂ fuel cell, 0.2 mgPt/cm², 80°C, 200 kPa back pressure")

add_image(sl, os.path.join(FIGS, "fig5.png"),
          Inches(0.3), Inches(1.15), width=Inches(7.8))

add_rect(sl, Inches(8.2), Inches(1.2), Inches(4.9), Inches(5.8), fill_color=NAVY)
txb(sl, "MEA Results Summary", Inches(8.3), Inches(1.3), Inches(4.6), Inches(0.45),
    font_size=16, bold=True, color=WHITE)

mea = [
    ("Metric",          "Pt₃Co/Te-C",   "DOE Target"),
    ("BOL MA @ 0.9V",   "0.50 A/mgPt",  "≥ 0.44"),
    ("EOL MA @ 0.9V",   "0.37 A/mgPt",  "—"),
    ("MA Retention",    "74% ✓",        "≥ 60%"),
    ("Power density",   "2.32 W/cm²\n@4 A/cm²", "—"),
    ("Pt₃Co/C BOL MA",  "0.165 A/mgPt  (4.3×\nlower than Te-C)", "—"),
]
for ri, row in enumerate(mea):
    bg = BLUE if ri==0 else (TEAL if ri in (1,3) else RGBColor(0x1E,0x3A,0x5F) if ri%2==0 else RGBColor(0x14,0x2A,0x44))
    y = Inches(1.9) + ri*Inches(0.75)
    add_rect(sl, Inches(8.3), y, Inches(4.6), Inches(0.7), fill_color=bg)
    txb(sl, f"{row[0]}: {row[1]}  |  {row[2]}",
        Inches(8.35), y+Inches(0.03), Inches(4.5), Inches(0.65),
        font_size=13, color=WHITE)

txb(sl, "Both BOL and EOL metrics significantly surpass DOE 2025.",
    Inches(8.3), Inches(6.3), Inches(4.6), Inches(0.45),
    font_size=14, bold=True, color=RGBColor(0x88,0xFF,0xCC))

# ─── SLIDE 14 – Post-durability MEA Analysis ────────────────────────────────
sl = blank_slide(prs)
fill_bg(sl, LGRAY)
header_bar(sl, "Post-Durability MEA: L1₂ Structure Preserved on Te-C",
           subtitle="Figures 5d–i — SAXS, HAADF-STEM at BOL and EOL")

add_image(sl, os.path.join(FIGS, "fig5.png"),
          Inches(0.3), Inches(1.15), width=Inches(7.8))

bullet_box(sl, [
    "SAXS analysis: BOL NP size ~4 nm → EOL predominantly ~3 nm",
    "HAADF-STEM EOL: average size 3.01 nm — consistent with SAXS",
    "Atomic-resolution STEM at EOL: L1₂ ordered structure intact on Te-C",
    "Pt₃Co/C at EOL: L1₂ structure LOST — large Co dissolution (see Table S8)",
    "Te-C strongly retains Co within i-NPs, preventing disordering",
    "Size decrease (4→3 nm) attributed to Pt dissolution during Ostwald ripening",
    "Despite some dissolution, structural integrity preserved → Te is doing its job",
], Inches(8.2), Inches(1.3), Inches(4.9), Inches(5.5), font_size=16, color=DKGRAY)

# ─── SLIDE 15 – DFT Calculations ────────────────────────────────────────────
sl = blank_slide(prs)
fill_bg(sl, LGRAY)
header_bar(sl, "DFT Calculations: Mechanistic Understanding of MSI and ORR",
           subtitle="Figure 6 — Binding energy, charge density, ORR free energy pathway")

add_image(sl, os.path.join(FIGS, "fig6.png"),
          Inches(0.3), Inches(1.15), width=Inches(7.8))

add_rect(sl, Inches(8.2), Inches(1.2), Inches(4.9), Inches(5.8), fill_color=NAVY)
txb(sl, "DFT Findings", Inches(8.3), Inches(1.3), Inches(4.7), Inches(0.45),
    font_size=16, bold=True, color=WHITE)

dft = [
    ("Binding Energies (Pt₁₀Co₃ cluster on graphene):", ""),
    ("  Te-C:  −10.28 eV  (strongest)", ""),
    ("  Se-C:  −9.49 eV", ""),
    ("  S-C:   −8.61 eV", ""),
    ("  C:     −3.95 eV  (weakest)", ""),
    ("Charge density: Te doping increases", "electron transfer metal → support"),
    ("ORR rate-limiting step: *OH desorption", ""),
    ("Limiting potentials:", ""),
    ("  Te-C@Pt(111): 0.64 V  ← best", ""),
    ("  Se-C: 0.62 V | S-C: 0.61 V | C: 0.58 V", ""),
]
tb = slide_shapes_add_textbox = sl.shapes.add_textbox(Inches(8.3), Inches(1.9), Inches(4.6), Inches(5.0))
tf = tb.text_frame
tf.word_wrap = True
first = True
for label, _ in dft:
    if first:
        p = tf.paragraphs[0]; first = False
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(3)
    run = p.add_run()
    run.text = label
    run.font.size = Pt(14)
    run.font.color.rgb = WHITE if "strongest" not in label and "best" not in label else RGBColor(0x88,0xFF,0xCC)
    run.font.bold = "strongest" in label or "best" in label

# ─── SLIDE 16 – Conclusions ──────────────────────────────────────────────────
sl = blank_slide(prs)
fill_bg(sl, NAVY)
add_rect(sl, 0, 0, Inches(0.18), H, fill_color=TEAL)

txb(sl, "Conclusions", Inches(0.35), Inches(0.3), Inches(12.5), Inches(0.7),
    font_size=34, bold=True, color=WHITE)
add_rect(sl, Inches(0.35), Inches(1.05), Inches(4.0), Inches(0.06), fill_color=TEAL)

conclusions = [
    "Te-modified carbon (Te-C) synthesised by CVD anchors Pt₃Co i-NPs with "
    "strong Pt–Te bonds, greatly enhancing metal–support interactions (MSI)",
    "Smaller NPs (3.9 nm vs 9.1 nm) due to suppressed sintering → higher ECSA and mass activity",
    "RDE: only 1.5% SA loss after 100,000 cycles — morphology essentially unchanged",
    "MEA: BOL MA 0.50 A/mgPt; EOL retention 74% — both exceed DOE 2025 requirements",
    "DFT confirms Te-C gives strongest binding energy (−10.28 eV) and optimal *OH desorption",
    "Thermostability of Te is the key: 2.3 wt% remains after annealing vs. ~0 for S",
    "This Te-anchoring strategy opens a new avenue for designing durable Pt-based i-NP catalysts",
]
bullet_box(sl, conclusions, Inches(0.35), Inches(1.2), Inches(12.4), Inches(5.8),
           font_size=18, color=WHITE, bullet="✦  ")

# ─── SLIDE 17 – Take-home / Discussion ──────────────────────────────────────
sl = blank_slide(prs)
fill_bg(sl, LGRAY)
header_bar(sl, "Discussion & Outlook")

left_items = [
    "What makes this work stand out?",
    "   • Te anchoring is a simple, scalable CVD process",
    "   • 100k RDE cycles without morphological change is record-breaking",
    "   • Clear mechanistic picture from XAS + DFT (Pt–Te bond confirmed)",
    "",
    "Remaining questions:",
    "   • Long-term real-device stability (>30k MEA cycles)?",
    "   • Cost of Te at scale vs. performance gain?",
    "   • Applicability to other intermetallic systems (PtNi, PtFe)?",
    "   • Behaviour under H₂–Air (vs H₂–O₂ tested here)?",
]
bullet_box(sl, left_items, Inches(0.3), Inches(1.3), Inches(6.2), Inches(5.8),
           font_size=16, color=DKGRAY, bullet="")

add_rect(sl, Inches(6.7), Inches(1.3), Inches(6.3), Inches(5.8), fill_color=NAVY)
txb(sl, "Key Takeaway", Inches(6.8), Inches(1.4), Inches(6.1), Inches(0.5),
    font_size=18, bold=True, color=TEAL)
txb(sl,
    "Tellurium's unique thermostability — surviving high-temperature "
    "synthesis conditions that destroy S and Se — is the decisive factor "
    "that enables this catalyst to simultaneously achieve small NP size, "
    "strong metal–support interaction, ordered intermetallic structure, "
    "and record-breaking durability.\n\n"
    "This is not just an incremental improvement; it represents a "
    "design principle: choose a heteroatom dopant that persists under "
    "the conditions your catalyst will experience.",
    Inches(6.8), Inches(2.0), Inches(6.1), Inches(5.0),
    font_size=16, color=WHITE)

# ── save ──────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Total slides: {len(prs.slides)}")
